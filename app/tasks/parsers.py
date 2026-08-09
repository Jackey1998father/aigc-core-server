"""
文档解析器（按文件类型分发）
"""
import csv
import logging
from io import BytesIO, StringIO
from typing import Callable, Dict

logger = logging.getLogger(__name__)


# ==================== 解析器实现 ====================

def _parse_pdf(file_bytes: bytes) -> str:
    """
    PDF → Markdown（保留标题、表格、列表）
    使用 pymupdf + pymupdf4llm，输出格式直接适合 RAG。
    """
    import pymupdf
    import pymupdf4llm
    doc = pymupdf.open(stream=file_bytes, filetype="pdf")
    return pymupdf4llm.to_markdown(doc)


def _parse_docx(file_bytes: bytes) -> str:
    """
    Word → 段落文本（也提取表格里的内容）
    """
    from docx import Document

    doc = Document(BytesIO(file_bytes))
    texts = []

    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                texts.append(row_text)

    return "\n\n".join(texts)


def _parse_pptx(file_bytes: bytes) -> str:
    """
    PPT → 每页幻灯片的文本（带页码标记）
    """
    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    slides_md = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"# 第{i}页"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_lines.append(para.text)
        if len(slide_lines) > 1:  # 有实际内容
            slides_md.append("\n\n".join(slide_lines))

    return "\n\n---\n\n".join(slides_md)


def _parse_xlsx(file_bytes: bytes) -> str:
    """
    Excel → 每个 sheet 转为 Markdown 表格
    餐饮场景：菜品成本表、供应商列表、采购价格表等。
    """
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(file_bytes), data_only=True)
    sheets_md = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        sheets_md.append(f"## Sheet: {sheet_name}")

        # 第一行作为表头
        header = [str(c) if c is not None else "" for c in rows[0]]
        sheets_md.append("| " + " | ".join(header) + " |")
        sheets_md.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            row_str = [str(c) if c is not None else "" for c in row]
            sheets_md.append("| " + " | ".join(row_str) + " |")

    return "\n\n".join(sheets_md)


def _parse_csv(file_bytes: bytes) -> str:
    """
    CSV → Markdown 表格
    """
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(StringIO(text))
    rows = list(reader)
    if not rows:
        return ""

    lines = []
    header = rows[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        # 补齐列数
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")

    return "\n".join(lines)


def _parse_plain_text(file_bytes: bytes) -> str:
    """
    txt/md → 直接读取
    """
    return file_bytes.decode("utf-8", errors="ignore")


# ==================== 分发入口 ====================

# 文件类型 → 解析器 的映射表（统一在这里维护）
PARSERS: Dict[str, Callable[[bytes], str]] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "xlsx": _parse_xlsx,
    "xls": _parse_xlsx,
    "csv": _parse_csv,
    "txt": _parse_plain_text,
    "md": _parse_plain_text,
}


def parse_by_type(file_type: str, file_bytes: bytes) -> str:
    """
    按文件类型分发到对应解析器，返回 Markdown 文本。

    Args:
        file_type: 文件扩展名（pdf/docx/pptx/xlsx/csv/txt/md...）
        file_bytes: 文件字节流

    Returns:
        解析后的文本（Markdown 格式）
    """
    parser = PARSERS.get((file_type or "").lower())
    if parser is None:
        logger.warning("[parsers] 不支持的文件类型: %s，按纯文本尝试", file_type)
        return _parse_plain_text(file_bytes)

    return parser(file_bytes)